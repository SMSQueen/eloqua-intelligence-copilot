# app.py
# Eloqua Intelligence Copilot — Polished UI (executive demo)
# - Branded header + sidebar controls
# - KPIs, charts (matplotlib), styled tables
# - One-click exports: Markdown, PDF, PPTX
# - Relative paths + synthetic data fallback

import streamlit as st
import pandas as pd
import numpy as np
from pathlib import Path
from datetime import date, timedelta
import io
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.pagesizes import LETTER
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib.pyplot as plt

st.set_page_config(page_title="Eloqua Intelligence Copilot", page_icon="📊", layout="wide")

# ---------- Paths & data safety ----------
REPO_ROOT = Path(__file__).resolve().parent
DATA_DIR = REPO_ROOT / "ingest" / "example_data"
DATA_DIR.mkdir(parents=True, exist_ok=True)

def ensure_example_data():
    emails_csv = DATA_DIR / "emails_by_segment_daily.csv"
    segwin_csv = DATA_DIR / "segment_engagement_windows.csv"
    if emails_csv.exists() and segwin_csv.exists():
        return
    today = date(2025, 9, 28)
    dates = [today - timedelta(days=i) for i in range(27, -1, -1)]
    segments = [
        {"segment_id": "MPE", "segment_name": "Owner Relations – East"},
        {"segment_id": "MPW", "segment_name": "Owner Relations – West"},
        {"segment_id": "CORP", "segment_name": "Corporate HQ"},
    ]
    rows = []
    rng = np.random.default_rng(42)
    for dt in dates:
        for s in segments:
            sends = rng.integers(600, 1400)
            opens = int(sends * rng.uniform(0.25, 0.40))
            clicks = int(opens * rng.uniform(0.07, 0.13))
            unsubs = rng.integers(0, max(1, sends // 1500))
            rows.append({
                "date": dt.isoformat(),
                "segment_id": s["segment_id"],
                "segment_name": s["segment_name"],
                "sends": int(sends),
                "opens": int(opens),
                "clicks": int(clicks),
                "unsubs": int(unsubs),
                "spam_complaints": 0,
                "unique_contacts_reached": int(sends * rng.uniform(0.7, 0.95)),
            })
    df = pd.DataFrame(rows)
    df.to_csv(emails_csv, index=False)
    dfr = df.copy()
    dfr["date"] = pd.to_datetime(dfr["date"])
    dfr = dfr.sort_values(["segment_id", "date"])
    dfr["EPC_7d"] = dfr.groupby("segment_id")["sends"].transform(lambda s: s.rolling(7, 1).sum()) / \
                    dfr.groupby("segment_id")["unique_contacts_reached"].transform(lambda s: s.rolling(7, 1).sum()).replace(0, np.nan)
    dfr["ctor_7d"] = dfr.groupby("segment_id")["clicks"].transform(lambda s: s.rolling(7, 1).sum()) / \
                     dfr.groupby("segment_id")["opens"].transform(lambda s: s.rolling(7, 1).sum()).replace(0, np.nan)
    dfr["unsub_rate_7d"] = dfr.groupby("segment_id")["unsubs"].transform(lambda s: s.rolling(7, 1).sum()) / \
                           dfr.groupby("segment_id")["sends"].transform(lambda s: s.rolling(7, 1).sum()).replace(0, np.nan)
    dfr["oversaturation_flag"] = (dfr["EPC_7d"] > 4)
    dfr["fatigue_score"] = (dfr["EPC_7d"].fillna(0) / 5.0).clip(0, 1)
    dfr.to_csv(segwin_csv, index=False)

ensure_example_data()
df = pd.read_csv(DATA_DIR / "emails_by_segment_daily.csv", parse_dates=["date"])
df_roll = pd.read_csv(DATA_DIR / "segment_engagement_windows.csv", parse_dates=["date"])

# ---------- Brand header ----------
with st.container():
    left, right = st.columns([0.7, 0.3])
    with left:
        st.markdown(
            "<h1 style='margin-bottom:0'>Eloqua Intelligence Copilot</h1>"
            "<div style='color:#00467F;margin-top:2px'><b>Executive Insights • Fatigue Optimizer • A/B Test Guidance</b></div>",
            unsafe_allow_html=True
        )
    with right:
        st.markdown(
            "<div style='text-align:right;'>"
            "<span style='background:#E6F0FA;color:#00467F;padding:6px 10px;border-radius:8px;'>Portfolio Demo</span>"
            "</div>",
            unsafe_allow_html=True
        )
st.write("")

# ---------- Sidebar filters ----------
with st.sidebar:
    st.markdown("### Controls")
    end_date = df["date"].max().date()
    start_date = (df["date"].max() - pd.Timedelta(days=6)).date()
    start = st.date_input("Start date", value=start_date, min_value=df["date"].min().date(), max_value=end_date)
    end = st.date_input("End date", value=end_date, min_value=start, max_value=df["date"].max().date())
    segments = sorted(df["segment_name"].unique().tolist())
    selected_segs = st.multiselect("Segments", options=segments, default=segments)
    st.markdown("---")
    st.caption("Tip: Narrow to 1–2 segments to see clearer trends.")

mask = (df["date"].dt.date >= start) & (df["date"].dt.date <= end)
if selected_segs:
    mask &= df["segment_name"].isin(selected_segs)
week = df[mask]
last_week_mask = (df["date"].dt.date >= (start - timedelta(days=(end - start).days + 1))) & (df["date"].dt.date < start)
last_week = df[last_week_mask & df["segment_name"].isin(selected_segs)]

# ---------- KPI cards ----------
def kpi(label, current, previous=None, pct=False):
    delta = None
    if previous is not None and previous != 0:
        delta = (current - previous) / previous * 100.0
    st.metric(label, f"{current:.2%}" if pct else f"{int(current)}",
              None if delta is None else f"{delta:+.1f}%")

c1, c2, c3, c4 = st.columns(4)
with c1: kpi("Sends", week["sends"].sum(), last_week["sends"].sum() if not last_week.empty else None, pct=False)
with c2: kpi("Open Rate", week["opens"].sum()/max(1,week["sends"].sum()),
             last_week["opens"].sum()/max(1,last_week["sends"].sum()) if not last_week.empty else None, pct=True)
with c3: kpi("CTR", week["clicks"].sum()/max(1,week["sends"].sum()),
             last_week["clicks"].sum()/max(1,last_week["sends"].sum()) if not last_week.empty else None, pct=True)
with c4: kpi("CTOR", week["clicks"].sum()/max(1,week["opens"].sum()),
             last_week["clicks"].sum()/max(1,last_week["opens"].sum()) if not last_week.empty else None, pct=True)

st.divider()

# ---------- Segment table ----------
st.subheader("Segment Performance")
seg = week.groupby("segment_name").agg(
    sends=("sends","sum"),
    opens=("opens","sum"),
    clicks=("clicks","sum"),
    unsubs=("unsubs","sum")
)
if not seg.empty:
    seg["open_rate"] = seg["opens"]/seg["sends"]
    seg["ctr"] = seg["clicks"]/seg["sends"]
    seg["ctor"] = seg["clicks"]/seg["opens"].clip(lower=1)
    st.dataframe(seg.reset_index(), use_container_width=True)
else:
    st.info("No data in the selected window/segments.")

# ---------- Charts (matplotlib, one per figure, default colors) ----------
chart_cols = st.columns(2)
if not seg.empty:
    # CTOR by segment (bar)
    with chart_cols[0]:
        fig1 = plt.figure()
        plt.bar(seg.index, seg["ctor"])
        plt.title("CTOR by Segment")
        plt.xticks(rotation=20, ha="right")
        plt.ylabel("CTOR")
        st.pyplot(fig1)

# 7-day CTOR trend (line) for selected segments
trend = df[df["segment_name"].isin(selected_segs)].copy()
trend["date"] = pd.to_datetime(trend["date"])
trend = trend.groupby(["segment_name","date"]).agg(opens=("opens","sum"), clicks=("clicks","sum")).reset_index()
trend["ctor"] = trend["clicks"]/trend["opens"].replace(0, np.nan)
trend = trend[(trend["date"].dt.date >= start) & (trend["date"].dt.date <= end)]

with chart_cols[1]:
    if not trend.empty:
        fig2 = plt.figure()
        for name, g in trend.groupby("segment_name"):
            plt.plot(g["date"], g["ctor"], marker="o", label=name)
        plt.title("CTOR Trend (Selected Segments)")
        plt.xlabel("Date"); plt.ylabel("CTOR")
        plt.xticks(rotation=20, ha="right")
        plt.legend()
        st.pyplot(fig2)

st.divider()

# ---------- Fatigue alerts ----------
st.subheader("Fatigue & Frequency Alerts")
latest_day = df_roll["date"].max()
latest = df_roll[df_roll["date"] == latest_day]
latest = latest[latest["segment_id"].isin(df[df["segment_name"].isin(selected_segs)]["segment_id"].unique())]
alerts = latest[(latest.get("oversaturation_flag", False)) | (latest.get("fatigue_score", 0) > 0.6)]
if alerts.empty:
    st.success("No fatigue alerts for the latest day. 👍")
else:
    show_cols = [c for c in ["date","segment_id","EPC_7d","ctor_7d","unsub_rate_7d","fatigue_score","oversaturation_flag"] if c in alerts.columns]
    st.dataframe(alerts[show_cols].reset_index(drop=True), use_container_width=True)

st.divider()

# ---------- Executive brief + exports ----------
st.subheader("Executive-Ready Weekly Brief")

def make_brief_md(week_df, last_week_df, seg_df):
    if week_df.empty:
        return "# Brief\nNo data in selected window."
    top_seg = seg_df["ctor"].idxmax() if not seg_df.empty else "N/A"
    lag_seg = seg_df["ctor"].idxmin() if not seg_df.empty else "N/A"
    this_ctor = week_df["clicks"].sum()/max(1,week_df["opens"].sum())
    last_ctor = last_week_df["clicks"].sum()/max(1,last_week_df["opens"].sum()) if not last_week_df.empty else 0
    delta_ctor = this_ctor - last_ctor
    md = f"""
# Eloqua Performance Brief ({start} to {end})

**Headlines**
- Overall CTOR: {this_ctor:.2%} (Δ vs prior window: {delta_ctor:+.2%})
- Top segment: **{top_seg}** (CTOR {seg_df.loc[top_seg,"ctor"]:.2% if not seg_df.empty else 0:.2%})
- Lagging segment: **{lag_seg}** (CTOR {seg_df.loc[lag_seg,"ctor"]:.2% if not seg_df.empty else 0:.2%})

**Fatigue & Risk**
- Review latest-day alerts above. Throttle where EPC-7d > 4 and CTOR-7d breaks below baseline.

**Next Best Tests**
- Subject: Verb-led, 5–7 words vs urgency phrase
- CTA: “Review & Confirm” vs “Confirm Now”
- Send-time: 09:30 vs 13:00
"""
    return md

brief_md = make_brief_md(week, last_week, seg if not seg.empty else pd.DataFrame())
st.code(brief_md, language="markdown")

colA, colB, colC = st.columns(3)

with colA:
    st.download_button("⬇️ Download Brief (.md)", brief_md.encode("utf-8"), file_name="weekly_brief.md", mime="text/markdown")

with colB:
    # Generate PDF on the fly
    def build_pdf(md_text: str) -> bytes:
        styles = getSampleStyleSheet()
        story = []
        for line in md_text.splitlines():
            if line.startswith("# "):
                story.append(Paragraph(f"<b><font size=16>{line[2:].strip()}</font></b>", styles["Normal"]))
                story.append(Spacer(1, 12))
            elif line.startswith("**") and line.endswith("**"):
                story.append(Paragraph(f"<b>{line.strip('*')}</b>", styles["Normal"]))
            elif line.startswith("**"):
                story.append(Spacer(1, 8))
                story.append(Paragraph(f"<b>{line.strip('*')}</b>", styles["Normal"]))
                story.append(Spacer(1, 4))
            elif line.strip().startswith("- "):
                story.append(Paragraph("• " + line.strip()[2:], styles["Normal"]))
            elif line.strip() == "":
                story.append(Spacer(1, 6))
            else:
                story.append(Paragraph(line, styles["Normal"]))
        buff = io.BytesIO()
        SimpleDocTemplate(buff, pagesize=LETTER).build(story)
        return buff.getvalue()

    pdf_bytes = build_pdf(brief_md)
    st.download_button("⬇️ Download Brief (PDF)", pdf_bytes, file_name="weekly_brief.pdf", mime="application/pdf")

with colC:
    # Generate PPTX on the fly
    def build_pptx(md_text: str) -> bytes:
        prs = Presentation()
        prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
        # Title
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(2))
        p = tb.text_frame.add_paragraph(); p.text = "Eloqua Performance Brief"; p.font.size = Pt(48); p.font.bold = True; p.font.color.rgb = RGBColor(0,70,127); p.alignment = PP_ALIGN.CENTER
        sub = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.33), Inches(1))
        sp = sub.text_frame.add_paragraph(); sp.text = f"{start} – {end}"; sp.font.size = Pt(24); sp.font.color.rgb = RGBColor(90,90,90); sp.alignment = PP_ALIGN.CENTER
        # Sections from MD
        lines = md_text.splitlines()
        section, bullets = None, []
        def add_bullets(title, items):
            s = prs.slides.add_slide(prs.slide_layouts[6])
            t = s.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(1))
            tp = t.text_frame.add_paragraph(); tp.text = title; tp.font.size = Pt(32); tp.font.bold = True; tp.font.color.rgb = RGBColor(0,70,127)
            body = s.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(5)).text_frame
            for b in items:
                p = body.add_paragraph(); p.text = b; p.font.size = Pt(20)
        for ln in lines:
            if ln.startswith("**") and ln.endswith("**"):
                if section and bullets: add_bullets(section, bullets); bullets=[]
                section = ln.strip("*")
            elif ln.strip().startswith("- "):
                bullets.append(ln.strip()[2:])
        if section and bullets: add_bullets(section, bullets)
        bio = io.BytesIO(); prs.save(bio); return bio.getvalue()

    pptx_bytes = build_pptx(brief_md)
    st.download_button("⬇️ Download Brief (PPTX)", pptx_bytes, file_name="weekly_brief.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
